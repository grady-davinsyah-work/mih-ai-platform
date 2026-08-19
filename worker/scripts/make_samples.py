"""Generate sample documents for smoke testing."""
import argparse
from pathlib import Path

from pptx import Presentation
from docx import Document
from reportlab.pdfgen import canvas


def write_pptx(path: Path):
    prs = Presentation()
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Paparan Rencana Pembangunan — Slide {i + 1}"
        body = slide.placeholders[1].text_frame
        body.text = ("Kedeputian Perencanaan Makro merancang arah pembangunan "
                     "jangka menengah. Prioritas meliputi pertumbuhan ekonomi "
                     "inklusif dan pemerataan infrastruktur.")
        slide.notes_slide.notes_text_frame.text = f"Catatan pembicara slide {i + 1}."
    prs.save(str(path))


def write_pdf(path: Path):
    c = canvas.Canvas(str(path))
    for i in range(2):
        c.drawString(72, 740, f"Laporan Perencanaan Makro — Halaman {i + 1}")
        c.drawString(72, 720, "Dokumen ini memuat indikator makro pembangunan nasional.")
        c.showPage()
    c.save()


def write_docx(path: Path):
    doc = Document()
    doc.add_heading("Pendahuluan", level=1)
    doc.add_paragraph("Laporan ini menjelaskan kondisi makro ekonomi terkini.")
    doc.add_heading("Arah Kebijakan", level=1)
    doc.add_paragraph("Kebijakan difokuskan pada stabilitas dan pertumbuhan.")
    doc.save(str(path))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("outdir")
    args = ap.parse_args()
    out = Path(args.outdir)
    out.mkdir(parents=True, exist_ok=True)
    write_pptx(out / "paparan-rencana.pptx")
    write_pdf(out / "laporan-makro.pdf")
    write_docx(out / "laporan-makro.docx")
    print("sampel dibuat di", out)


if __name__ == "__main__":
    main()

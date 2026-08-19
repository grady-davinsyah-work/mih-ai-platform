from pathlib import Path

from .base import Segment


def _first_line(text: str) -> str | None:
    line = next((ln.strip() for ln in text.splitlines() if ln.strip()), None)
    return line or None


def parse_pptx(path: Path) -> list[Segment]:
    from pptx import Presentation

    prs = Presentation(str(path))
    segments: list[Segment] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
            elif shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Catatan pembicara] {notes}")
        text = "\n".join(texts).strip()
        if text:
            segments.append(
                Segment(text=text, page_or_slide=i, section_title=_first_line(text))
            )
    return segments

from pptx import Presentation
from docx import Document
from reportlab.pdfgen import canvas

from parsers import parse_pptx, parse_pdf, parse_docx


def make_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Judul Slide 1"
    body = slide.placeholders[1].text_frame
    body.text = "Isi paparan slide pertama."
    slide.notes_slide.notes_text_frame.text = "Catatan pembicara di sini."
    prs.save(path)


def make_pdf(path):
    c = canvas.Canvas(str(path))
    c.drawString(72, 700, "Halaman satu dari laporan.")
    c.showPage()
    c.drawString(72, 700, "Halaman dua dari laporan.")
    c.save()


def make_pdf_blank(path):
    c = canvas.Canvas(str(path))
    c.showPage()
    c.save()


def make_docx(path):
    doc = Document()
    doc.add_heading("Pendahuluan", level=1)
    doc.add_paragraph("Paragraf pembuka laporan.")
    doc.add_heading("Metodologi", level=1)
    doc.add_paragraph("Penjelasan metodologi.")
    doc.save(path)


def test_parse_pptx_slides_and_notes(tmp_path):
    p = tmp_path / "d.pptx"
    make_pptx(p)
    segs = parse_pptx(p)
    assert len(segs) == 1
    assert segs[0].page_or_slide == 1
    assert "Isi paparan" in segs[0].text
    assert "Catatan pembicara" in segs[0].text


def test_parse_pdf_text_per_page(tmp_path):
    p = tmp_path / "a.pdf"
    make_pdf(p)
    segs = parse_pdf(p)
    assert len(segs) == 2
    assert segs[0].page_or_slide == 1
    assert "Halaman satu" in segs[0].text
    assert segs[0].needs_ocr is False


def test_parse_pdf_blank_flags_ocr(tmp_path):
    p = tmp_path / "b.pdf"
    make_pdf_blank(p)
    segs = parse_pdf(p)
    assert segs[0].needs_ocr is True


def test_parse_docx_grouped_by_heading(tmp_path):
    p = tmp_path / "c.docx"
    make_docx(p)
    segs = parse_docx(p)
    assert len(segs) == 2
    assert segs[0].section_title == "Pendahuluan"
    assert "Paragraf pembuka" in segs[0].text
    assert segs[1].section_title == "Metodologi"
